import io
import os
import shutil
import subprocess
import tempfile
import zipfile
import base64
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Optional

import streamlit as st
import streamlit.components.v1 as components
from pypdf import PdfReader, PdfWriter


APP_TITLE = "PDF & Image Tools"

# Streamlit multipage can execute the main script as `__main__`, while tool pages
# do `import app` / `from app import *`. That can result in two copies of this
# file being loaded and the `app` import becoming stale after a rerun.
# Alias the currently-running module as `app` so pages always import the latest.
if __name__ == "__main__":
    sys.modules["app"] = sys.modules[__name__]


@dataclass
class DependencyStatus:
    poppler: bool
    tesseract: bool
    libreoffice: bool
    ghostscript: bool
    qpdf: bool

    # Resolved binary locations (when found). Useful on macOS where Streamlit's PATH
    # can differ from the interactive shell.
    pdftoppm_cmd: Optional[str] = None
    poppler_path: Optional[str] = None  # directory containing pdftoppm
    tesseract_cmd: Optional[str] = None
    soffice_cmd: Optional[str] = None
    gs_cmd: Optional[str] = None
    qpdf_cmd: Optional[str] = None


def _is_executable_file(path: str | Path) -> bool:
    try:
        p = Path(path)
        return p.is_file() and os.access(str(p), os.X_OK)
    except Exception:
        return False


def which(cmd: str, *, extra_dirs: Iterable[str] = ()) -> Optional[str]:
    found = shutil.which(cmd)
    if found:
        return found

    path = os.environ.get("PATH", "")
    for d in extra_dirs:
        if not d:
            continue
        found = shutil.which(cmd, path=str(d) + os.pathsep + path)
        if found:
            return found
    return None


def resolve_executable(cmd: str, *, candidates: Iterable[str] = ()) -> Optional[str]:
    # 1) Normal PATH
    found = which(cmd)
    if found:
        return found

    # 2) Common macOS install locations (Homebrew Intel + Apple Silicon)
    common_dirs = [
        "/opt/local/bin",
        "/opt/local/sbin",
        "/opt/homebrew/bin",
        "/opt/homebrew/sbin",
        "/usr/local/bin",
        "/usr/local/sbin",
        "/usr/bin",
        "/bin",
    ]
    found = which(cmd, extra_dirs=common_dirs)
    if found:
        return found

    # 3) Tool-specific candidates (including .app bundles)
    for c in candidates:
        if c and _is_executable_file(c):
            return str(Path(c))

    return None


def get_dependency_status() -> DependencyStatus:
    # macOS common binaries:
    # - poppler: pdftoppm
    # - tesseract: tesseract
    # - libreoffice: soffice
    # - ghostscript: gs
    pdftoppm_cmd = resolve_executable(
        "pdftoppm",
        candidates=[
            "/opt/local/bin/pdftoppm",
            "/opt/homebrew/opt/poppler/bin/pdftoppm",
            "/usr/local/opt/poppler/bin/pdftoppm",
            "/opt/homebrew/bin/pdftoppm",
            "/usr/local/bin/pdftoppm",
        ],
    )
    tesseract_cmd = resolve_executable(
        "tesseract",
        candidates=[
            "/opt/local/bin/tesseract",
            "/opt/homebrew/bin/tesseract",
            "/usr/local/bin/tesseract",
            "/opt/homebrew/opt/tesseract/bin/tesseract",
            "/usr/local/opt/tesseract/bin/tesseract",
        ],
    )
    lo_candidates: list[str] = [
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice.bin",
    ]
    try:
        apps = Path("/Applications")
        for p in apps.glob("LibreOffice*.app/Contents/MacOS/soffice"):
            lo_candidates.append(str(p))
        for p in apps.glob("LibreOffice*.app/Contents/MacOS/soffice.bin"):
            lo_candidates.append(str(p))
    except Exception:
        pass

    soffice_cmd = resolve_executable("soffice", candidates=lo_candidates)
    gs_cmd = resolve_executable(
        "gs",
        candidates=[
            "/opt/local/bin/gs",
            "/opt/homebrew/bin/gs",
            "/usr/local/bin/gs",
            "/opt/homebrew/opt/ghostscript/bin/gs",
            "/usr/local/opt/ghostscript/bin/gs",
        ],
    )
    qpdf_cmd = resolve_executable(
        "qpdf",
        candidates=[
            "/opt/local/bin/qpdf",
            "/opt/homebrew/bin/qpdf",
            "/usr/local/bin/qpdf",
            "/opt/homebrew/opt/qpdf/bin/qpdf",
            "/usr/local/opt/qpdf/bin/qpdf",
        ],
    )

    poppler_path = str(Path(pdftoppm_cmd).parent) if pdftoppm_cmd else None

    return DependencyStatus(
        poppler=pdftoppm_cmd is not None,
        tesseract=tesseract_cmd is not None,
        libreoffice=soffice_cmd is not None,
        ghostscript=gs_cmd is not None,
        qpdf=qpdf_cmd is not None,
        pdftoppm_cmd=pdftoppm_cmd,
        poppler_path=poppler_path,
        tesseract_cmd=tesseract_cmd,
        soffice_cmd=soffice_cmd,
        gs_cmd=gs_cmd,
        qpdf_cmd=qpdf_cmd,
    )


def pdf2image_convert_from_bytes(
    pdf_bytes: bytes,
    *,
    fmt: str,
    poppler_path: Optional[str] = None,
    dpi: Optional[int] = None,
    first_page: Optional[int] = None,
    last_page: Optional[int] = None,
):
    from pdf2image import convert_from_bytes

    kwargs: dict = {"fmt": fmt}
    if dpi is not None:
        kwargs["dpi"] = int(dpi)
    if first_page is not None:
        kwargs["first_page"] = int(first_page)
    if last_page is not None:
        kwargs["last_page"] = int(last_page)
    if poppler_path:
        kwargs["poppler_path"] = poppler_path
    return convert_from_bytes(pdf_bytes, **kwargs)


def tesseract_list_langs(*, tesseract_cmd: str = "tesseract") -> list[str]:
    # Returns language codes like: eng, tam, osd...
    try:
        out = run_cmd([tesseract_cmd, "--list-langs"])
        text = out.decode("utf-8", errors="replace")
        langs: list[str] = []
        for line in text.splitlines():
            line = line.strip()
            if not line or line.lower().startswith("list of available"):
                continue
            langs.append(line)
        return sorted(set(langs))
    except Exception:
        return []


def _uploaded_file_bytes(uploaded_file) -> bytes:
    if uploaded_file is None:
        return b""
    getvalue = getattr(uploaded_file, "getvalue", None)
    if callable(getvalue):
        val = getvalue()
        if isinstance(val, (bytes, bytearray, memoryview)):
            return bytes(val)
        return b""
    read = getattr(uploaded_file, "read", None)
    if callable(read):
        data = read()
        seek = getattr(uploaded_file, "seek", None)
        if callable(seek):
            try:
                seek(0)
            except Exception:
                pass
        if isinstance(data, (bytes, bytearray, memoryview)):
            return bytes(data)
        return b""
    return b""


def uploaded_file_bytes(uploaded_file) -> bytes:
    return _uploaded_file_bytes(uploaded_file)


def _is_pdf_upload(uploaded_file) -> bool:
    if uploaded_file is None:
        return False
    name = (getattr(uploaded_file, "name", "") or "").lower()
    if name.endswith(".pdf"):
        return True
    mime = (getattr(uploaded_file, "type", "") or "").lower()
    return mime == "application/pdf" or mime.endswith("/pdf")


def _is_image_upload(uploaded_file) -> bool:
    if uploaded_file is None:
        return False
    mime = (getattr(uploaded_file, "type", "") or "").lower()
    if mime.startswith("image/"):
        return True
    name = (getattr(uploaded_file, "name", "") or "").lower()
    return any(name.endswith(ext) for ext in [".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp", ".tif", ".tiff", ".heic", ".heif"])


def _preview_pdf_bytes(pdf_bytes: bytes, *, deps: DependencyStatus, file_name: str = "document.pdf"):
    # Render first page as an image (avoids Chrome blocking embedded PDFs).
    # Prefer PyMuPDF (pure Python wheels) and fall back to Poppler/pdf2image when available.

    # 1) PyMuPDF
    try:
        import fitz  # type: ignore

        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        if doc.page_count >= 1:
            page = doc.load_page(0)
            pix = page.get_pixmap(dpi=140)
            png_bytes = pix.tobytes("png")
            st.image(png_bytes, caption=f"Preview: {file_name} (page 1)", use_container_width=True)
            return
    except Exception:
        pass

    # 2) Poppler/pdf2image
    if deps.poppler:
        try:
            images = pdf2image_convert_from_bytes(
                pdf_bytes,
                fmt="png",
                dpi=140,
                first_page=1,
                last_page=1,
                poppler_path=deps.poppler_path,
            )
            if images:
                st.image(images[0], caption=f"Preview: {file_name} (page 1)", use_container_width=True)
                return
        except Exception:
            pass

    st.info("Preview unavailable. Install `pymupdf` (recommended) or Poppler (`pdftoppm`) to preview PDFs.")


def _preview_uploaded(uploaded, *, deps: DependencyStatus):
    if uploaded is None:
        return

    # Multiple file uploader returns a list.
    if isinstance(uploaded, list):
        if not uploaded:
            return
        # Preview all PDFs (page 1) as requested; also preview images.
        any_previewed = False
        for item in uploaded:
            if _is_pdf_upload(item) or _is_image_upload(item):
                _preview_uploaded(item, deps=deps)
                any_previewed = True
        if not any_previewed:
            first = uploaded[0]
            st.caption(f"Previewing first upload: {getattr(first, 'name', '')}")
            return _preview_uploaded(first, deps=deps)
        return

    if _is_image_upload(uploaded):
        try:
            img = pil_open_image(uploaded)
            st.image(img, caption=f"Preview: {getattr(uploaded, 'name', 'image')}", use_container_width=True)
        except Exception as e:
            st.caption(f"Preview unavailable: {e}")
        return

    if _is_pdf_upload(uploaded):
        pdf_bytes = uploaded_file_bytes(uploaded)
        if not pdf_bytes:
            return
        _preview_pdf_bytes(pdf_bytes, deps=deps, file_name=getattr(uploaded, "name", "document.pdf"))
        return


def _install_file_uploader_preview_patch():
    # Patch Streamlit's file_uploader once so every tool page gets previews.
    if getattr(st, "_pdf_tools_uploader_patched", False):
        return
    st._pdf_tools_uploader_patched = True  # type: ignore[attr-defined]

    orig = st.file_uploader

    def wrapped_file_uploader(*args, **kwargs):
        uploaded = orig(*args, **kwargs)
        try:
            deps = get_dependency_status()
            _preview_uploaded(uploaded, deps=deps)
        except Exception:
            # Never block tool flows due to preview issues.
            pass
        return uploaded

    st.file_uploader = wrapped_file_uploader  # type: ignore[assignment]


_install_file_uploader_preview_patch()


def run_cmd(args: list[str], *, input_bytes: Optional[bytes] = None) -> bytes:
    proc = subprocess.run(
        args,
        input=input_bytes,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        check=False,
    )
    if proc.returncode != 0:
        raise RuntimeError(proc.stderr.decode("utf-8", errors="replace")[:4000])
    return proc.stdout


def rewrite_pdf_with_pypdf(pdf_bytes: bytes) -> bytes:
    reader = PdfReader(io.BytesIO(pdf_bytes))
    writer = PdfWriter()
    if reader.is_encrypted:
        raise ValueError("PDF is encrypted; unlock it first.")
    for page in reader.pages:
        writer.add_page(page)
    return writer_to_bytes(writer)


def qpdf_optimize(pdf_bytes: bytes, *, qpdf_cmd: str = "qpdf") -> bytes:
    # qpdf is the most reliable open-source optimizer on macOS.
    with tempfile.TemporaryDirectory() as td:
        in_path = Path(td) / "in.pdf"
        out_path = Path(td) / "out.pdf"
        in_path.write_bytes(pdf_bytes)
        run_cmd([qpdf_cmd, "--linearize", "--stream-data=compress", str(in_path), str(out_path)])
        return out_path.read_bytes()


def read_uploaded_pdf(file) -> PdfReader:
    data = file.read()
    return PdfReader(io.BytesIO(data))


def writer_to_bytes(writer: PdfWriter) -> bytes:
    out = io.BytesIO()
    writer.write(out)
    return out.getvalue()


def download_button(label: str, data: bytes, filename: str, mime: str = "application/pdf"):
    st.download_button(label=label, data=data, file_name=filename, mime=mime)


def zip_files(files: list[tuple[str, bytes]]) -> bytes:
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, mode="w", compression=zipfile.ZIP_DEFLATED) as zf:
        for name, data in files:
            zf.writestr(name, data)
    return buf.getvalue()


def sort_uploads(files: Optional[list]) -> list:
    if not files:
        return []
    return sorted(list(files), key=lambda f: (getattr(f, "name", "") or "").casefold())


def show_sorted_names(files: Optional[list], *, label: str = "Processing order (A→Z)") -> list:
    ordered = sort_uploads(files)
    if ordered:
        names = [getattr(f, "name", "") for f in ordered]
        st.caption(f"{label}: " + ", ".join(names))
    return ordered


def pil_open_image(uploaded_file):
    from PIL import Image

    # Optional HEIC/HEIF support.
    try:
        import pillow_heif  # type: ignore

        pillow_heif.register_heif_opener()
    except Exception:
        pass

    uploaded_file.seek(0)
    img = Image.open(uploaded_file)
    # Keep alpha when present; convert palette images.
    if img.mode in ("P", "LA"):
        img = img.convert("RGBA")
    return img


def pil_to_bytes(img, fmt: str, *, quality: int = 85) -> bytes:
    buf = io.BytesIO()
    fmt_u = fmt.upper()
    save_kwargs: dict = {}
    if fmt_u in ("JPG", "JPEG"):
        if img.mode not in ("RGB",):
            img = img.convert("RGB")
        save_kwargs.update({"quality": int(quality), "optimize": True, "progressive": True})
        fmt_u = "JPEG"
    elif fmt_u == "PNG":
        save_kwargs.update({"optimize": True})
    elif fmt_u == "WEBP":
        save_kwargs.update({"quality": int(quality), "method": 6})
    img.save(buf, format=fmt_u, **save_kwargs)
    return buf.getvalue()


def simple_remove_background(img, *, threshold: int = 245):
    # Simple offline background remover for near-white backgrounds.
    from PIL import Image

    rgba = img.convert("RGBA")
    pixels = rgba.load()
    w, h = rgba.size
    for y in range(h):
        for x in range(w):
            r, g, b, a = pixels[x, y]
            if r >= threshold and g >= threshold and b >= threshold:
                pixels[x, y] = (r, g, b, 0)
    return rgba


def blur_boxes(img, boxes_px: list[tuple[int, int, int, int]], *, radius: int = 12):
    from PIL import ImageFilter

    out = img.convert("RGB")
    for (x1, y1, x2, y2) in boxes_px:
        left, right = sorted([x1, x2])
        top, bottom = sorted([y1, y2])
        crop = out.crop((left, top, right, bottom)).filter(ImageFilter.GaussianBlur(radius=radius))
        out.paste(crop, (left, top, right, bottom))
    return out


def extract_pdf_text(reader: PdfReader, *, max_chars: int = 2_000_000) -> str:
    # Best-effort text extraction; capped to avoid huge memory usage.
    chunks: list[str] = []
    size = 0
    for page in reader.pages:
        t = page.extract_text() or ""
        if not t:
            continue
        if size + len(t) > max_chars:
            remaining = max_chars - size
            if remaining > 0:
                chunks.append(t[:remaining])
            break
        chunks.append(t)
        size += len(t)
    return "\n\n".join(chunks)


def simple_summarize(text: str, *, sentence_count: int = 6) -> str:
    import re

    cleaned = " ".join((text or "").split())
    if not cleaned:
        return ""

    # Very small, offline summarizer: sentence scoring by word frequency.
    sentences = re.split(r"(?<=[.!?])\s+", cleaned)
    sentences = [s.strip() for s in sentences if len(s.strip()) > 20]
    if not sentences:
        return cleaned[:1000]

    words = re.findall(r"[A-Za-z']{2,}", cleaned.lower())
    stop = {
        "the",
        "and",
        "a",
        "an",
        "to",
        "of",
        "in",
        "is",
        "it",
        "for",
        "on",
        "with",
        "as",
        "by",
        "at",
        "from",
        "that",
        "this",
        "be",
        "are",
        "was",
        "were",
        "or",
        "but",
        "not",
        "we",
        "you",
        "i",
        "they",
        "he",
        "she",
        "them",
        "his",
        "her",
        "their",
        "our",
    }
    freq: dict[str, int] = {}
    for w in words:
        if w in stop:
            continue
        freq[w] = freq.get(w, 0) + 1

    def score_sentence(s: str) -> float:
        sw = re.findall(r"[A-Za-z']{2,}", s.lower())
        if not sw:
            return 0.0
        return sum(freq.get(w, 0) for w in sw) / (len(sw) ** 0.5)

    ranked = sorted(
        ((i, score_sentence(s), s) for i, s in enumerate(sentences)),
        key=lambda x: x[1],
        reverse=True,
    )
    top = sorted(ranked[: max(1, min(sentence_count, len(sentences)))], key=lambda x: x[0])
    return " ".join(s for _, __, s in top)


def text_to_pdf_bytes(text: str, *, title: str = "Document") -> bytes:
    # Simple text-to-PDF renderer; keeps things dependency-light.
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=letter)
    width, height = letter
    c.setFont("Helvetica-Bold", 14)
    c.drawString(40, height - 40, title)
    c.setFont("Helvetica", 10)

    y = height - 70
    for raw_line in (text or "").splitlines() or [""]:
        line = raw_line.rstrip()
        # crude wrap
        while len(line) > 110:
            c.drawString(40, y, line[:110])
            line = line[110:]
            y -= 14
            if y < 40:
                c.showPage()
                c.setFont("Helvetica", 10)
                y = height - 40
        c.drawString(40, y, line)
        y -= 14
        if y < 40:
            c.showPage()
            c.setFont("Helvetica", 10)
            y = height - 40

    c.save()
    return buf.getvalue()


def reorder_pages(reader: PdfReader, new_order: list[int]) -> bytes:
    total = len(reader.pages)
    if not new_order:
        raise ValueError("New order is empty")
    if any(p < 1 or p > total for p in new_order):
        raise ValueError("New order contains out-of-range page numbers")
    writer = PdfWriter()
    for p in new_order:
        writer.add_page(reader.pages[p - 1])
    return writer_to_bytes(writer)


def parse_page_list(text: str) -> list[int]:
    # Accept: "3,1,2,4" (reorder list)
    text = (text or "").strip()
    if not text:
        return []
    return [int(x.strip()) for x in text.split(",") if x.strip()]


def redact_via_rasterize(
    pdf_bytes: bytes,
    boxes: list[tuple[int, float, float, float, float]],
    *,
    poppler_path: Optional[str] = None,
) -> bytes:
    # boxes: (page, x1, y1, x2, y2) normalized 0..1, origin top-left.
    from PIL import ImageDraw

    images = pdf2image_convert_from_bytes(pdf_bytes, fmt="png", dpi=200, poppler_path=poppler_path)
    by_page: dict[int, list[tuple[float, float, float, float]]] = {}
    for page, x1, y1, x2, y2 in boxes:
        by_page.setdefault(page, []).append((x1, y1, x2, y2))

    out_images = []
    for idx, img in enumerate(images, start=1):
        draw = ImageDraw.Draw(img)
        w, h = img.size
        for (x1, y1, x2, y2) in by_page.get(idx, []):
            xa = int(max(0.0, min(1.0, x1)) * w)
            ya = int(max(0.0, min(1.0, y1)) * h)
            xb = int(max(0.0, min(1.0, x2)) * w)
            yb = int(max(0.0, min(1.0, y2)) * h)
            left, right = sorted([xa, xb])
            top, bottom = sorted([ya, yb])
            draw.rectangle([left, top, right, bottom], fill="black")
        out_images.append(img.convert("RGB"))

    buf = io.BytesIO()
    out_images[0].save(buf, format="PDF", save_all=True, append_images=out_images[1:])
    return buf.getvalue()


def stamp_signature(reader: PdfReader, signature_png: bytes, pages: set[int], x: float, y: float, w: float) -> bytes:
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader

    writer = PdfWriter()
    img_reader = ImageReader(io.BytesIO(signature_png))

    for idx, page in enumerate(reader.pages, start=1):
        if idx in pages:
            packet = io.BytesIO()
            c = canvas.Canvas(packet, pagesize=(page.mediabox.width, page.mediabox.height))
            c.drawImage(img_reader, x, y, width=w, preserveAspectRatio=True, mask="auto")
            c.save()
            packet.seek(0)
            overlay = PdfReader(packet)
            page.merge_page(overlay.pages[0])
        writer.add_page(page)
    return writer_to_bytes(writer)


def merge_pdfs(files: list) -> bytes:
    writer = PdfWriter()
    for f in sort_uploads(files):
        reader = read_uploaded_pdf(f)
        for page in reader.pages:
            writer.add_page(page)
    return writer_to_bytes(writer)


def split_pdf_by_ranges(reader: PdfReader, ranges: list[tuple[int, int]]) -> list[tuple[str, bytes]]:
    # ranges are 1-based inclusive: (start, end)
    outputs: list[tuple[str, bytes]] = []
    total = len(reader.pages)
    for idx, (start, end) in enumerate(ranges, start=1):
        start0 = max(1, start)
        end0 = min(total, end)
        if start0 > end0:
            continue
        writer = PdfWriter()
        for i in range(start0 - 1, end0):
            writer.add_page(reader.pages[i])
        outputs.append((f"part-{idx}-{start0}-{end0}.pdf", writer_to_bytes(writer)))
    return outputs


def parse_page_ranges(text: str) -> list[tuple[int, int]]:
    # Accept: "1-3, 5, 7-9"
    text = (text or "").strip()
    if not text:
        return []
    ranges: list[tuple[int, int]] = []
    for chunk in text.split(","):
        chunk = chunk.strip()
        if not chunk:
            continue
        if "-" in chunk:
            a, b = chunk.split("-", 1)
            ranges.append((int(a.strip()), int(b.strip())))
        else:
            n = int(chunk)
            ranges.append((n, n))
    return ranges


def remove_pages(reader: PdfReader, pages_to_remove: set[int]) -> bytes:
    # pages_to_remove is 1-based
    writer = PdfWriter()
    for idx, page in enumerate(reader.pages, start=1):
        if idx in pages_to_remove:
            continue
        writer.add_page(page)
    return writer_to_bytes(writer)


def rotate_pages(reader: PdfReader, pages: set[int], degrees: int) -> bytes:
    writer = PdfWriter()
    for idx, page in enumerate(reader.pages, start=1):
        if idx in pages:
            page.rotate(degrees)
        writer.add_page(page)
    return writer_to_bytes(writer)


def main():
    st.set_page_config(page_title=APP_TITLE, layout="wide")
    st.title(APP_TITLE)

    deps = get_dependency_status()

    with st.expander("Optional dependencies status", expanded=False):
        st.write(
            {
                "poppler (pdftoppm)": deps.poppler,
                "tesseract": deps.tesseract,
                "libreoffice (soffice)": deps.libreoffice,
                "ghostscript (gs)": deps.ghostscript,
                "qpdf": deps.qpdf,
            }
        )
        st.caption(
            "Resolved: "
            + ", ".join(
                [
                    f"pdftoppm={deps.pdftoppm_cmd or 'not found'}",
                    f"tesseract={deps.tesseract_cmd or 'not found'}",
                    f"soffice={deps.soffice_cmd or 'not found'}",
                    f"gs={deps.gs_cmd or 'not found'}",
                    f"qpdf={deps.qpdf_cmd or 'not found'}",
                ]
            )
        )
        st.caption(
            "Missing dependencies will only disable the related tools; the rest still works."
        )

        missing: list[str] = []
        if not deps.poppler:
            missing.append("poppler")
        if not deps.tesseract:
            missing.append("tesseract")
        if not deps.libreoffice:
            missing.append("libreoffice")
        if not deps.ghostscript:
            missing.append("ghostscript")
        if not deps.qpdf:
            missing.append("qpdf")

        if missing:
            is_macports = resolve_executable("port") is not None or (deps.tesseract_cmd or "").startswith("/opt/local/")
            if is_macports:
                st.caption("Install missing tools (MacPorts):")
                st.code("sudo port install " + " ".join(sorted(set(missing))), language="bash")
            else:
                st.caption("Install missing tools (Homebrew):")
                st.code("brew install " + " ".join(sorted(set(missing))), language="bash")

    # Home router: send the user to a dedicated Streamlit page per tool.
    pages_dir = Path(__file__).parent / "pages"
    page_files = sorted([p for p in pages_dir.glob("*.py") if not p.name.startswith("__")]) if pages_dir.exists() else []
    if page_files:
        tool_map: dict[str, str] = {}
        for p in page_files:
            stem = p.stem
            parts = stem.split("_", 1)
            if len(parts) == 2 and parts[0].isdigit():
                display = parts[1]
            else:
                display = stem
            display = display.replace("_", " ")
            tool_map[display] = f"pages/{p.name}"

        st.subheader("Open a tool")
        choice = st.selectbox("Tool", options=["— Select —", *tool_map.keys()], index=0)
        if choice and choice != "— Select —":
            st.switch_page(tool_map[choice])
            st.stop()
        st.caption("Tip: you can also use Streamlit's sidebar Pages list.")
        st.stop()

    # Layout inspired by the screenshot: categories and tools.
    col1, col2, col3, col4, col5 = st.columns(5)

    # ORGANIZE PDF
    with col1:
        st.subheader("Organize PDF")

        with st.expander("Merge PDF"):
            files = st.file_uploader(
                "Select PDFs", type=["pdf"], accept_multiple_files=True, key="merge"
            )
            ordered_files = show_sorted_names(files)
            if files and st.button("Merge", key="merge_btn"):
                out = merge_pdfs(ordered_files)
                download_button("Download merged PDF", out, "merged.pdf")

        with st.expander("Split PDF"):
            f = st.file_uploader("PDF", type=["pdf"], key="split")
            ranges_text = st.text_input(
                "Page ranges (e.g., 1-3, 5, 7-9)", key="split_ranges"
            )
            if f and st.button("Split", key="split_btn"):
                reader = read_uploaded_pdf(f)
                ranges = parse_page_ranges(ranges_text)
                parts = split_pdf_by_ranges(reader, ranges)
                if not parts:
                    st.warning("No valid ranges provided.")
                for name, data in parts:
                    download_button(f"Download {name}", data, name)

        with st.expander("Remove pages"):
            f = st.file_uploader("PDF", type=["pdf"], key="remove")
            pages_text = st.text_input("Pages to remove (e.g., 2, 4-6)", key="remove_pages")
            if f and st.button("Remove selected pages", key="remove_btn"):
                reader = read_uploaded_pdf(f)
                ranges = parse_page_ranges(pages_text)
                pages = set()
                for a, b in ranges:
                    pages.update(range(min(a, b), max(a, b) + 1))
                out = remove_pages(reader, pages)
                download_button("Download PDF", out, "pages-removed.pdf")

        with st.expander("Extract pages"):
            f = st.file_uploader("PDF", type=["pdf"], key="extract")
            ranges_text = st.text_input("Pages to extract (e.g., 1-2, 5)", key="extract_ranges")
            if f and st.button("Extract", key="extract_btn"):
                reader = read_uploaded_pdf(f)
                ranges = parse_page_ranges(ranges_text)
                parts = split_pdf_by_ranges(reader, ranges)
                if not parts:
                    st.warning("No valid page selection provided.")
                for name, data in parts:
                    download_button(f"Download {name}", data, name)

        with st.expander("Organize PDF"):
            st.info("Reorder pages by providing a new page order list.")
            f = st.file_uploader("PDF", type=["pdf"], key="organize")
            order_text = st.text_input("New order (e.g., 3,1,2,4)", key="organize_order")
            if f and st.button("Reorder", key="organize_btn"):
                reader = read_uploaded_pdf(f)
                try:
                    new_order = parse_page_list(order_text)
                    out = reorder_pages(reader, new_order)
                except Exception as e:
                    st.error(f"Reorder failed: {e}")
                    st.stop()
                download_button("Download reordered PDF", out, "reordered.pdf")

        with st.expander("Scan to PDF"):
            st.info("Upload images (photos/scans) and export as a single PDF.")
            imgs = st.file_uploader(
                "Select images",
                type=["jpg", "jpeg", "png", "webp", "heic", "heif", "tif", "tiff", "bmp", "gif"],
                accept_multiple_files=True,
                key="scan2pdf",
            )
            ordered_imgs = show_sorted_names(imgs)
            if imgs and st.button("Create scanned PDF", key="scan2pdf_btn"):
                from PIL import Image

                pil_imgs = [pil_open_image(f).convert("RGB") for f in ordered_imgs]
                out = io.BytesIO()
                pil_imgs[0].save(out, format="PDF", save_all=True, append_images=pil_imgs[1:])
                download_button("Download PDF", out.getvalue(), "scan.pdf")

        st.subheader("PDF Intelligence")

        with st.expander("AI Summarizer"):
            f = st.file_uploader("PDF", type=["pdf"], key="summarize")
            n = st.slider("Sentences", min_value=3, max_value=15, value=6)
            if f and st.button("Summarize", key="summarize_btn"):
                reader = read_uploaded_pdf(f)
                text = extract_pdf_text(reader)
                summary = simple_summarize(text, sentence_count=n)
                st.text_area("Summary", value=summary, height=200)
                st.download_button("Download summary (txt)", summary.encode("utf-8"), "summary.txt", "text/plain")
                st.download_button("Download summary (pdf)", text_to_pdf_bytes(summary, title="Summary"), "summary.pdf")

        with st.expander("Translate PDF"):
            st.info("Offline-first translation using Argos Translate if installed (no quotas).")
            f = st.file_uploader("PDF", type=["pdf"], key="translate")
            src = st.text_input("Source language code", value="en", key="tr_src")
            dst = st.text_input("Target language code", value="fr", key="tr_dst")
            has_argos = False
            argos_translate = None
            try:
                import argostranslate.package  # type: ignore
                import argostranslate.translate as _argos_translate  # type: ignore

                has_argos = True
                argos_translate = _argos_translate
            except Exception:
                has_argos = False

            if not has_argos:
                st.warning("Argos Translate not installed. Install with: pip install argostranslate")
            if f and st.button("Translate", key="translate_btn"):
                if not has_argos:
                    st.stop()
                reader = read_uploaded_pdf(f)
                text = extract_pdf_text(reader)
                translated = argos_translate.translate(text, from_code=src, to_code=dst)  # type: ignore[union-attr]
                st.text_area("Translated text", value=translated, height=200)
                st.download_button(
                    "Download translated (txt)",
                    translated.encode("utf-8"),
                    file_name=f"translated-{src}-to-{dst}.txt",
                    mime="text/plain",
                )
                st.download_button(
                    "Download translated (pdf)",
                    text_to_pdf_bytes(translated, title=f"Translated ({src}→{dst})"),
                    file_name=f"translated-{src}-to-{dst}.pdf",
                )

    # OPTIMIZE PDF
    with col2:
        st.subheader("Optimize PDF")
        st.caption("(Best results require optional system tools.)")

        with st.expander("Compress PDF"):
            st.info(
                "Runs a best-effort optimization. If `qpdf` is installed, it will linearize and compress streams; otherwise it will re-write the PDF via pypdf (lighter optimization)."
            )
            f = st.file_uploader("PDF", type=["pdf"], key="compress")
            if f and st.button("Compress/Optimize", key="compress_btn"):
                data = f.read()
                try:
                    out = (
                        qpdf_optimize(data, qpdf_cmd=deps.qpdf_cmd or "qpdf")
                        if deps.qpdf
                        else rewrite_pdf_with_pypdf(data)
                    )
                except Exception as e:
                    st.error(f"Optimization failed: {e}")
                    st.stop()
                download_button("Download optimized PDF", out, "optimized.pdf")

        with st.expander("Repair PDF"):
            st.info(
                "Attempts to open and re-save the PDF to fix minor structure issues. If `qpdf` is installed, it will also attempt a clean rewrite."
            )
            f = st.file_uploader("PDF", type=["pdf"], key="repair")
            if f and st.button("Repair", key="repair_btn"):
                data = uploaded_file_bytes(f)
                try:
                    if deps.qpdf:
                        out = qpdf_optimize(data, qpdf_cmd=deps.qpdf_cmd or "qpdf")
                    else:
                        out = rewrite_pdf_with_pypdf(data)
                except Exception as e:
                    st.error(f"Repair failed: {e}")
                    st.stop()
                download_button("Download repaired PDF", out, "repaired.pdf")

        with st.expander("OCR PDF"):
            st.info("Uses Tesseract + Poppler to produce a searchable PDF (best-effort).")
            if not deps.tesseract or not deps.poppler:
                st.warning("Requires 'tesseract' and 'pdftoppm' (poppler).")

            if deps.tesseract and "tesseract_langs" not in st.session_state:
                st.session_state["tesseract_langs"] = tesseract_list_langs(
                    tesseract_cmd=deps.tesseract_cmd or "tesseract"
                )
            available_langs = st.session_state.get("tesseract_langs", [])

            if available_langs:
                default_langs = [l for l in ["eng", "tam"] if l in available_langs]
                if not default_langs and "eng" in available_langs:
                    default_langs = ["eng"]
                langs = st.multiselect(
                    "OCR languages",
                    options=available_langs,
                    default=default_langs,
                    help="Tesseract language codes. Install extra language data on the host to add more options.",
                )
                if "tam" not in available_langs:
                    st.caption("Tamil (tam) is not installed on this host.")
            else:
                langs = st.multiselect(
                    "OCR languages",
                    options=["eng", "tam"],
                    default=["eng"],
                    help="Language list unavailable. This host may be missing Tesseract language data.",
                )

            lang_arg = "+".join([l for l in langs if l]) if langs else None
            ocr_config = "--oem 1 --psm 3 -c preserve_interword_spaces=1"
            st.caption(
                "OCR accuracy depends on scan quality; exact extraction cannot be guaranteed for all PDFs. "
                "For best results, use correct language(s) and a clean scan."
            )

            f = st.file_uploader("PDF", type=["pdf"], key="ocr")

            run_col, extract_col = st.columns(2)
            run_ocr = run_col.button("Run OCR", key="ocr_btn")
            extract_text = extract_col.button("Extract text", key="ocr_extract_btn")

            if f and (run_ocr or extract_text):
                if not deps.tesseract or not deps.poppler:
                    st.stop()
                import pytesseract

                pdf_bytes = uploaded_file_bytes(f)
                images = pdf2image_convert_from_bytes(
                    pdf_bytes, fmt="png", dpi=220, poppler_path=deps.poppler_path
                )

                if run_ocr:
                    writer = PdfWriter()
                    for img in images:
                        pdf_page_bytes = pytesseract.image_to_pdf_or_hocr(
                            img,
                            extension="pdf",
                            lang=lang_arg,
                            config=ocr_config,
                        )
                        if isinstance(pdf_page_bytes, str):
                            pdf_page_bytes = pdf_page_bytes.encode("utf-8")
                        page_reader = PdfReader(io.BytesIO(pdf_page_bytes))
                        writer.add_page(page_reader.pages[0])
                    download_button("Download searchable PDF", writer_to_bytes(writer), "ocr.pdf")

                if extract_text:
                    chunks: list[str] = []
                    for img in images:
                        txt = pytesseract.image_to_string(
                            img,
                            lang=lang_arg,
                            config=ocr_config,
                        )
                        chunks.append(txt or "")
                    extracted = "\n\n".join(chunks).strip()
                    st.text_area("Extracted text", value=extracted, height=420)
                    st.download_button(
                        "Download extracted text (txt)",
                        extracted.encode("utf-8"),
                        file_name="ocr-extracted.txt",
                        mime="text/plain",
                    )

    # CONVERT TO PDF
    with col3:
        st.subheader("Convert to PDF")

        with st.expander("JPG to PDF"):
            imgs = st.file_uploader(
                "Select images",
                type=["jpg", "jpeg", "png", "webp", "heic", "heif", "tif", "tiff", "bmp", "gif"],
                accept_multiple_files=True,
                key="img2pdf",
            )
            ordered_imgs = show_sorted_names(imgs)
            if imgs and st.button("Convert to PDF", key="img2pdf_btn"):
                from PIL import Image

                pil_imgs = []
                for f in ordered_imgs:
                    img = pil_open_image(f).convert("RGB")
                    pil_imgs.append(img)
                out = io.BytesIO()
                pil_imgs[0].save(out, format="PDF", save_all=True, append_images=pil_imgs[1:])
                download_button("Download PDF", out.getvalue(), "images.pdf")

        with st.expander("WORD to PDF"):
            st.info("Uses LibreOffice (soffice) if installed.")
            if not deps.libreoffice:
                st.warning("Requires 'soffice' (LibreOffice).")
            f = st.file_uploader("DOC/DOCX", type=["doc", "docx"], key="doc2pdf")
            if f and st.button("Convert", key="doc2pdf_btn"):
                if not deps.libreoffice:
                    st.stop()
                with tempfile.TemporaryDirectory() as td:
                    in_path = Path(td) / f.name
                    in_path.write_bytes(f.read())
                    subprocess.check_call([
                        deps.soffice_cmd or "soffice",
                        "--headless",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        td,
                        str(in_path),
                    ])
                    out_path = in_path.with_suffix(".pdf")
                    download_button("Download PDF", out_path.read_bytes(), out_path.name)

        with st.expander("POWERPOINT to PDF"):
            st.info("Uses LibreOffice (soffice) if installed.")
            if not deps.libreoffice:
                st.warning("Requires 'soffice' (LibreOffice).")
            f = st.file_uploader("PPT/PPTX", type=["ppt", "pptx"], key="ppt2pdf")
            if f and st.button("Convert", key="ppt2pdf_btn"):
                if not deps.libreoffice:
                    st.stop()
                with tempfile.TemporaryDirectory() as td:
                    in_path = Path(td) / f.name
                    in_path.write_bytes(f.read())
                    subprocess.check_call([
                        deps.soffice_cmd or "soffice",
                        "--headless",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        td,
                        str(in_path),
                    ])
                    out_path = in_path.with_suffix(".pdf")
                    download_button("Download PDF", out_path.read_bytes(), out_path.name)

        with st.expander("EXCEL to PDF"):
            st.info("Uses LibreOffice (soffice) if installed.")
            if not deps.libreoffice:
                st.warning("Requires 'soffice' (LibreOffice).")
            f = st.file_uploader("XLS/XLSX", type=["xls", "xlsx"], key="xls2pdf")
            if f and st.button("Convert", key="xls2pdf_btn"):
                if not deps.libreoffice:
                    st.stop()
                with tempfile.TemporaryDirectory() as td:
                    in_path = Path(td) / f.name
                    in_path.write_bytes(f.read())
                    subprocess.check_call([
                        deps.soffice_cmd or "soffice",
                        "--headless",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        td,
                        str(in_path),
                    ])
                    out_path = in_path.with_suffix(".pdf")
                    download_button("Download PDF", out_path.read_bytes(), out_path.name)

        with st.expander("HTML to PDF"):
            st.info("Uses WeasyPrint (pure Python + system libs).")
            html = st.text_area("HTML", height=200, key="html_input")
            if html and st.button("Convert", key="html2pdf_btn"):
                from weasyprint import HTML

                pdf_bytes = HTML(string=html).write_pdf()
                if not pdf_bytes:
                    st.error("HTML→PDF failed to produce output.")
                    st.stop()
                download_button("Download PDF", pdf_bytes, "document.pdf")

    # CONVERT FROM PDF
    with col4:
        st.subheader("Convert from PDF")

        with st.expander("PDF to JPG"):
            st.info("Uses Poppler via pdf2image.")
            if not deps.poppler:
                st.warning("Requires 'pdftoppm' (poppler).")
            f = st.file_uploader("PDF", type=["pdf"], key="pdf2jpg")
            if f and st.button("Convert", key="pdf2jpg_btn"):
                if not deps.poppler:
                    st.stop()

                images = pdf2image_convert_from_bytes(f.read(), fmt="jpeg", poppler_path=deps.poppler_path)
                for i, img in enumerate(images, start=1):
                    buf = io.BytesIO()
                    img.save(buf, format="JPEG", quality=95)
                    st.download_button(
                        f"Download page {i}",
                        buf.getvalue(),
                        file_name=f"page-{i}.jpg",
                        mime="image/jpeg",
                    )

        with st.expander("PDF to WORD"):
            st.info("Basic text extraction to DOCX (layout not preserved).")
            f = st.file_uploader("PDF", type=["pdf"], key="pdf2word")
            if f and st.button("Convert", key="pdf2word_btn"):
                from docx import Document

                reader = read_uploaded_pdf(f)
                doc = Document()
                for page in reader.pages:
                    text = page.extract_text() or ""
                    doc.add_paragraph(text)
                out = io.BytesIO()
                doc.save(out)
                st.download_button(
                    "Download DOCX",
                    out.getvalue(),
                    file_name="document.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )

        with st.expander("PDF to POWERPOINT"):
            st.info("Exports each page as an image slide (requires Poppler).")
            if not deps.poppler:
                st.warning("Requires 'pdftoppm' (poppler).")
            f = st.file_uploader("PDF", type=["pdf"], key="pdf2ppt")
            if f and st.button("Convert", key="pdf2ppt_btn"):
                if not deps.poppler:
                    st.stop()
                from pptx import Presentation
                from pptx.util import Inches

                images = pdf2image_convert_from_bytes(f.read(), fmt="png", poppler_path=deps.poppler_path)
                prs = Presentation()
                blank = prs.slide_layouts[6]
                for img in images:
                    slide = prs.slides.add_slide(blank)
                    img_buf = io.BytesIO()
                    img.save(img_buf, format="PNG")
                    img_buf.seek(0)
                    slide.shapes.add_picture(img_buf, Inches(0), Inches(0), width=prs.slide_width)
                out = io.BytesIO()
                prs.save(out)
                st.download_button(
                    "Download PPTX",
                    out.getvalue(),
                    file_name="slides.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )

        with st.expander("PDF to EXCEL"):
            st.info("Extracts tables is non-trivial; this tool exports page text into a spreadsheet (one sheet per page).")
            f = st.file_uploader("PDF", type=["pdf"], key="pdf2xls")
            if f and st.button("Convert", key="pdf2xls_btn"):
                import openpyxl

                reader = read_uploaded_pdf(f)
                wb = openpyxl.Workbook()
                for i, page in enumerate(reader.pages, start=1):
                    if i == 1 and wb.active is not None:
                        ws = wb.active
                        ws.title = "Page 1"
                    else:
                        ws = wb.create_sheet(title=f"Page {i}")
                    text = (page.extract_text() or "").splitlines()
                    for r, line in enumerate(text, start=1):
                        ws.cell(row=r, column=1, value=line)
                out = io.BytesIO()
                wb.save(out)
                st.download_button(
                    "Download XLSX",
                    out.getvalue(),
                    file_name="document.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )

        with st.expander("PDF to PDF/A"):
            st.info("Uses Ghostscript to convert to PDF/A (best-effort).")
            if not deps.ghostscript:
                st.warning("Requires 'gs' (Ghostscript).")
            f = st.file_uploader("PDF", type=["pdf"], key="pdf2pdfa")
            if f and st.button("Convert", key="pdf2pdfa_btn"):
                if not deps.ghostscript:
                    st.stop()
                with tempfile.TemporaryDirectory() as td:
                    in_path = Path(td) / "input.pdf"
                    out_path = Path(td) / "output_pdfa.pdf"
                    in_path.write_bytes(f.read())
                    subprocess.check_call(
                        [
                            deps.gs_cmd or "gs",
                            "-dPDFA",
                            "-dBATCH",
                            "-dNOPAUSE",
                            "-sProcessColorModel=DeviceRGB",
                            "-sDEVICE=pdfwrite",
                            "-sOutputFile=" + str(out_path),
                            str(in_path),
                        ]
                    )
                    download_button("Download PDF/A", out_path.read_bytes(), "document-pdfa.pdf")

    # EDIT PDF + SECURITY
    with col5:
        st.subheader("Edit & Security")

        with st.expander("Rotate PDF"):
            f = st.file_uploader("PDF", type=["pdf"], key="rotate")
            pages_text = st.text_input("Pages to rotate (e.g., 1-3, 5)", key="rotate_pages")
            degrees = st.selectbox("Degrees", options=[90, 180, 270], index=0)
            if f and st.button("Rotate", key="rotate_btn"):
                reader = read_uploaded_pdf(f)
                ranges = parse_page_ranges(pages_text)
                pages = set()
                for a, b in ranges:
                    pages.update(range(min(a, b), max(a, b) + 1))
                out = rotate_pages(reader, pages, degrees)
                download_button("Download rotated PDF", out, "rotated.pdf")

        with st.expander("Add page numbers"):
            st.info("Adds simple bottom-center page numbers.")
            f = st.file_uploader("PDF", type=["pdf"], key="pagenum")
            if f and st.button("Add numbers", key="pagenum_btn"):
                from reportlab.pdfgen import canvas
                from reportlab.lib.units import inch

                reader = read_uploaded_pdf(f)
                writer = PdfWriter()
                total = len(reader.pages)
                for i, page in enumerate(reader.pages, start=1):
                    packet = io.BytesIO()
                    c = canvas.Canvas(packet, pagesize=(page.mediabox.width, page.mediabox.height))
                    c.setFont("Helvetica", 10)
                    c.drawCentredString(float(page.mediabox.width) / 2, 0.5 * inch, f"{i} / {total}")
                    c.save()
                    packet.seek(0)
                    overlay = PdfReader(packet)
                    page.merge_page(overlay.pages[0])
                    writer.add_page(page)
                download_button("Download PDF", writer_to_bytes(writer), "page-numbered.pdf")

        with st.expander("Add watermark"):
            st.info("Adds a diagonal text watermark on all pages.")
            f = st.file_uploader("PDF", type=["pdf"], key="watermark")
            text = st.text_input("Watermark text", value="CONFIDENTIAL", key="watermark_text")
            if f and st.button("Apply watermark", key="watermark_btn"):
                from reportlab.pdfgen import canvas

                reader = read_uploaded_pdf(f)
                writer = PdfWriter()
                for page in reader.pages:
                    packet = io.BytesIO()
                    c = canvas.Canvas(packet, pagesize=(page.mediabox.width, page.mediabox.height))
                    c.saveState()
                    c.setFont("Helvetica-Bold", 48)
                    c.translate(float(page.mediabox.width) / 2, float(page.mediabox.height) / 2)
                    c.rotate(45)
                    c.drawCentredString(0, 0, text)
                    c.restoreState()
                    c.save()
                    packet.seek(0)
                    overlay = PdfReader(packet)
                    page.merge_page(overlay.pages[0])
                    writer.add_page(page)
                download_button("Download PDF", writer_to_bytes(writer), "watermarked.pdf")

        with st.expander("Crop PDF"):
            st.info("Sets a crop box for all pages.")
            f = st.file_uploader("PDF", type=["pdf"], key="crop")
            left = st.number_input("Left", min_value=0.0, value=0.0)
            bottom = st.number_input("Bottom", min_value=0.0, value=0.0)
            right = st.number_input("Right (from left)", min_value=1.0, value=500.0)
            top = st.number_input("Top (from bottom)", min_value=1.0, value=700.0)
            if f and st.button("Apply crop", key="crop_btn"):
                reader = read_uploaded_pdf(f)
                writer = PdfWriter()
                for page in reader.pages:
                    page.cropbox.lower_left = (left, bottom)
                    page.cropbox.upper_right = (left + right, bottom + top)
                    writer.add_page(page)
                download_button("Download cropped PDF", writer_to_bytes(writer), "cropped.pdf")

        with st.expander("Edit PDF"):
            st.info("Full PDF editing is broad. Current app supports rotate, crop, watermark, and page numbers.")

        with st.expander("Unlock PDF"):
            st.info("Removes a password if you provide it (no brute forcing).")
            f = st.file_uploader("Encrypted PDF", type=["pdf"], key="unlock")
            pwd = st.text_input("Password", type="password", key="unlock_pwd")
            if f and st.button("Unlock", key="unlock_btn"):
                data = f.read()
                reader = PdfReader(io.BytesIO(data))
                if reader.is_encrypted:
                    ok = reader.decrypt(pwd)
                    if not ok:
                        st.error("Wrong password.")
                        st.stop()
                writer = PdfWriter()
                for page in reader.pages:
                    writer.add_page(page)
                download_button("Download unlocked PDF", writer_to_bytes(writer), "unlocked.pdf")

        with st.expander("Protect PDF"):
            st.info("Encrypts the PDF with a password.")
            f = st.file_uploader("PDF", type=["pdf"], key="protect")
            pwd = st.text_input("New password", type="password", key="protect_pwd")
            if f and st.button("Protect", key="protect_btn"):
                if not pwd:
                    st.warning("Password is required.")
                    st.stop()
                reader = read_uploaded_pdf(f)
                writer = PdfWriter()
                for page in reader.pages:
                    writer.add_page(page)
                writer.encrypt(pwd)
                download_button("Download protected PDF", writer_to_bytes(writer), "protected.pdf")

        with st.expander("Sign PDF"):
            st.info("Adds a visible signature stamp image (not a cryptographic certificate signature).")
            f = st.file_uploader("PDF", type=["pdf"], key="sign_pdf")
            sig = st.file_uploader("Signature image (PNG)", type=["png"], key="sign_img")
            pages_text = st.text_input("Pages to sign (e.g., 1, 3-4)", key="sign_pages")
            x = st.number_input("X (points)", value=50.0)
            y = st.number_input("Y (points)", value=50.0)
            w = st.number_input("Width (points)", value=200.0, min_value=10.0)
            if f and sig and st.button("Apply signature", key="sign_btn"):
                reader = read_uploaded_pdf(f)
                ranges = parse_page_ranges(pages_text)
                pages = set()
                for a, b in ranges:
                    pages.update(range(min(a, b), max(a, b) + 1))
                try:
                    out = stamp_signature(reader, sig.read(), pages, x, y, w)
                except Exception as e:
                    st.error(f"Signing failed: {e}")
                    st.stop()
                download_button("Download signed PDF", out, "signed.pdf")

        with st.expander("Redact PDF"):
            st.info(
                "Safe redaction by rasterizing pages and burning black boxes into the pixels. Requires Poppler (`pdftoppm`)."
            )
            if not deps.poppler:
                st.warning("Requires 'pdftoppm' (poppler).")
            f = st.file_uploader("PDF", type=["pdf"], key="redact_pdf")
            st.caption("Boxes format (normalized coords, origin top-left): page,x1,y1,x2,y2 per line. Example: 1,0.1,0.2,0.9,0.3")
            boxes_text = st.text_area("Redaction boxes", height=120, key="redact_boxes")
            if f and st.button("Redact", key="redact_btn"):
                if not deps.poppler:
                    st.stop()
                redact_boxes: list[tuple[int, float, float, float, float]] = []
                for line in (boxes_text or "").splitlines():
                    line = line.strip()
                    if not line:
                        continue
                    parts = [p.strip() for p in line.split(",")]
                    if len(parts) != 5:
                        st.error(f"Invalid line: {line}")
                        st.stop()
                    page = int(parts[0])
                    x1, y1, x2, y2 = map(float, parts[1:])
                    redact_boxes.append((page, x1, y1, x2, y2))
                try:
                    out = redact_via_rasterize(f.read(), redact_boxes, poppler_path=deps.poppler_path)
                except Exception as e:
                    st.error(f"Redaction failed: {e}")
                    st.stop()
                download_button("Download redacted PDF", out, "redacted.pdf")

        with st.expander("Compare PDF"):
            st.info("Compares extracted text per page (quick diff). For pixel-perfect compare, we can add image-based compare.")
            a = st.file_uploader("PDF A", type=["pdf"], key="cmp_a")
            b = st.file_uploader("PDF B", type=["pdf"], key="cmp_b")
            if a and b and st.button("Compare", key="cmp_btn"):
                from rapidfuzz import fuzz

                ra = read_uploaded_pdf(a)
                rb = read_uploaded_pdf(b)
                pages = max(len(ra.pages), len(rb.pages))
                results = []
                for i in range(pages):
                    ta = ra.pages[i].extract_text() if i < len(ra.pages) else ""
                    tb = rb.pages[i].extract_text() if i < len(rb.pages) else ""
                    score = fuzz.ratio(ta or "", tb or "")
                    results.append({"page": i + 1, "similarity": score})
                st.dataframe(results, use_container_width=True)

    st.divider()

    # IMAGE TOOLS
    st.header("Image Tools")
    st.caption("Local processing with no in-app quotas; limited only by your machine resources.")

    img_col1, img_col2, img_col3, img_col4, img_col5 = st.columns(5)

    with img_col1:
        with st.container(border=True):
            st.subheader("Compress IMAGE")
            st.caption("Compress JPG/PNG/WebP by re-encoding.")
            files = st.file_uploader(
                "Images",
                type=["jpg", "jpeg", "png", "webp", "heic", "heif", "tif", "tiff", "bmp", "gif"],
                accept_multiple_files=True,
                key="img_compress",
            )
            ordered_files = show_sorted_names(files)
            quality = st.slider("Quality", 10, 95, 80, key="img_compress_q")
            out_fmt = st.selectbox("Output", ["jpeg", "webp", "png"], index=0, key="img_compress_fmt")
            if files and st.button("Compress", key="img_compress_btn"):
                outs: list[tuple[str, bytes]] = []
                for f in ordered_files:
                    img = pil_open_image(f)
                    data = pil_to_bytes(img, out_fmt, quality=quality)
                    name = Path(f.name).stem + "." + ("jpg" if out_fmt == "jpeg" else out_fmt)
                    outs.append((name, data))
                z = zip_files(outs) if len(outs) > 1 else outs[0][1]
                if len(outs) > 1:
                    st.download_button("Download ZIP", z, file_name="compressed-images.zip", mime="application/zip")
                else:
                    st.download_button(
                        "Download image",
                        z,
                        file_name=outs[0][0],
                        mime="image/jpeg" if out_fmt == "jpeg" else f"image/{out_fmt}",
                    )

        with st.container(border=True):
            st.subheader("Rotate IMAGE")
            st.caption("Rotate JPG/PNG by 90/180/270.")
            f = st.file_uploader(
                "Image", type=["jpg", "jpeg", "png", "webp", "heic", "heif", "tif", "tiff", "bmp", "gif"], key="img_rotate"
            )
            deg = st.selectbox("Degrees", [90, 180, 270], index=0, key="img_rotate_deg")
            if f and st.button("Rotate", key="img_rotate_btn"):
                img = pil_open_image(f)
                out = img.rotate(-deg, expand=True)
                data = pil_to_bytes(out, "png")
                st.download_button("Download", data, file_name=f"{Path(f.name).stem}-rotated.png", mime="image/png")

    with img_col2:
        with st.container(border=True):
            st.subheader("Resize IMAGE")
            st.caption("Resize by percent or exact width/height.")
            f = st.file_uploader(
                "Image", type=["jpg", "jpeg", "png", "webp", "heic", "heif", "tif", "tiff", "bmp", "gif"], key="img_resize"
            )
            mode = st.radio("Mode", ["Percent", "Exact"], horizontal=True, key="img_resize_mode")
            keep = st.checkbox("Keep aspect ratio", value=True, key="img_resize_keep")
            if f:
                img = pil_open_image(f)
                w0, h0 = img.size
                st.image(img, caption=f"{w0}×{h0}", use_container_width=True)
                if mode == "Percent":
                    pct = st.slider("Scale %", 10, 400, 100, key="img_resize_pct")
                    w = max(1, int(w0 * pct / 100))
                    h = max(1, int(h0 * pct / 100))
                else:
                    w = st.number_input("Width", min_value=1, value=w0, key="img_resize_w")
                    h = st.number_input("Height", min_value=1, value=h0, key="img_resize_h")
                    if keep:
                        # adjust height when width changes (simple)
                        h = max(1, int(h0 * (int(w) / w0)))

                if st.button("Resize", key="img_resize_btn"):
                    out = img.resize((int(w), int(h)))
                    data = pil_to_bytes(out, "png")
                    st.download_button("Download", data, file_name=f"{Path(f.name).stem}-resized.png", mime="image/png")

        with st.container(border=True):
            st.subheader("Upscale Image")
            st.caption("Upscale with high-quality resampling (offline).")
            f = st.file_uploader(
                "Image", type=["jpg", "jpeg", "png", "webp", "heic", "heif", "tif", "tiff", "bmp", "gif"], key="img_upscale"
            )
            factor = st.selectbox("Scale", [2, 3, 4], index=0, key="img_upscale_factor")
            if f and st.button("Upscale", key="img_upscale_btn"):
                from PIL import Image

                img = pil_open_image(f)
                w0, h0 = img.size
                out = img.resize((w0 * factor, h0 * factor), resample=Image.Resampling.LANCZOS)
                data = pil_to_bytes(out, "png")
                st.download_button("Download", data, file_name=f"{Path(f.name).stem}-upscaled.png", mime="image/png")

    with img_col3:
        with st.container(border=True):
            st.subheader("Crop IMAGE")
            st.caption("Crop by pixel coordinates.")
            f = st.file_uploader(
                "Image", type=["jpg", "jpeg", "png", "webp", "heic", "heif", "tif", "tiff", "bmp", "gif"], key="img_crop"
            )
            if f:
                img = pil_open_image(f)
                w0, h0 = img.size
                st.image(img, use_container_width=True)
                x1 = st.number_input("Left", min_value=0, value=0, key="img_crop_x1")
                y1 = st.number_input("Top", min_value=0, value=0, key="img_crop_y1")
                x2 = st.number_input("Right", min_value=1, value=w0, key="img_crop_x2")
                y2 = st.number_input("Bottom", min_value=1, value=h0, key="img_crop_y2")
                if st.button("Crop", key="img_crop_btn"):
                    out = img.crop((int(x1), int(y1), int(x2), int(y2)))
                    data = pil_to_bytes(out, "png")
                    st.download_button("Download", data, file_name=f"{Path(f.name).stem}-cropped.png", mime="image/png")

        with st.container(border=True):
            st.subheader("Remove background")
            st.caption("Best for near-white backgrounds (offline).")
            f = st.file_uploader(
                "Image", type=["png", "jpg", "jpeg", "webp", "heic", "heif", "tif", "tiff", "bmp", "gif"], key="img_rmbg"
            )
            thresh = st.slider("White threshold", 200, 255, 245, key="img_rmbg_t")
            if f and st.button("Remove BG", key="img_rmbg_btn"):
                img = pil_open_image(f)
                out = None
                try:
                    from rembg import remove  # type: ignore

                    out_bytes = remove(f.getvalue())
                    st.download_button("Download", out_bytes, file_name=f"{Path(f.name).stem}-nobg.png", mime="image/png")
                except Exception:
                    out = simple_remove_background(img, threshold=int(thresh))
                    data = pil_to_bytes(out, "png")
                    st.download_button("Download", data, file_name=f"{Path(f.name).stem}-nobg.png", mime="image/png")

    with img_col4:
        with st.container(border=True):
            st.subheader("Convert to JPG")
            st.caption("Turn PNG/WebP/TIFF/etc into JPG.")
            files = st.file_uploader(
                "Images",
                type=["png", "webp", "heic", "heif", "tif", "tiff", "bmp", "gif", "jpg", "jpeg"],
                accept_multiple_files=True,
                key="img_to_jpg",
            )
            ordered_files = show_sorted_names(files)
            quality = st.slider("JPG quality", 10, 95, 85, key="img_to_jpg_q")
            if files and st.button("Convert", key="img_to_jpg_btn"):
                outs = []
                for f in ordered_files:
                    img = pil_open_image(f)
                    data = pil_to_bytes(img, "jpeg", quality=quality)
                    outs.append((Path(f.name).stem + ".jpg", data))
                z = zip_files(outs) if len(outs) > 1 else outs[0][1]
                if len(outs) > 1:
                    st.download_button("Download ZIP", z, file_name="converted-jpg.zip", mime="application/zip")
                else:
                    st.download_button("Download JPG", z, file_name=outs[0][0], mime="image/jpeg")

        with st.container(border=True):
            st.subheader("HTML to IMAGE")
            st.caption("Renders HTML to PDF then converts pages to PNG (needs Poppler).")
            html = st.text_area("HTML", height=120, key="html_to_img")
            if not deps.poppler:
                st.warning("Requires 'pdftoppm' (poppler).")
            if html and st.button("Render", key="html_to_img_btn"):
                if not deps.poppler:
                    st.stop()
                from weasyprint import HTML

                pdf_bytes = HTML(string=html).write_pdf()
                if not pdf_bytes:
                    st.error("HTML→PDF failed to produce output.")
                    st.stop()
                images = pdf2image_convert_from_bytes(pdf_bytes, fmt="png", dpi=200, poppler_path=deps.poppler_path)
                outs = []
                for i, img in enumerate(images, start=1):
                    outs.append((f"page-{i}.png", pil_to_bytes(img, "png")))
                z = zip_files(outs) if len(outs) > 1 else outs[0][1]
                if len(outs) > 1:
                    st.download_button("Download ZIP", z, file_name="html-images.zip", mime="application/zip")
                else:
                    st.download_button("Download PNG", z, file_name=outs[0][0], mime="image/png")

    with img_col5:
        with st.container(border=True):
            st.subheader("Convert from JPG")
            st.caption("Convert JPG to PNG/WebP/PDF.")
            files = st.file_uploader("JPG images", type=["jpg", "jpeg"], accept_multiple_files=True, key="img_from_jpg")
            ordered_files = show_sorted_names(files)
            target = st.selectbox("Target", ["png", "webp", "pdf"], index=0, key="img_from_jpg_target")
            quality = st.slider("Quality (WebP)", 10, 95, 85, key="img_from_jpg_q")
            if files and st.button("Convert", key="img_from_jpg_btn"):
                if target == "pdf":
                    from PIL import Image

                    pil_imgs = [pil_open_image(f).convert("RGB") for f in ordered_files]
                    buf = io.BytesIO()
                    pil_imgs[0].save(buf, format="PDF", save_all=True, append_images=pil_imgs[1:])
                    st.download_button("Download PDF", buf.getvalue(), file_name="images.pdf", mime="application/pdf")
                else:
                    outs = []
                    for f in ordered_files:
                        img = pil_open_image(f)
                        data = pil_to_bytes(img, target, quality=quality)
                        outs.append((Path(f.name).stem + f".{target}", data))
                    z = zip_files(outs) if len(outs) > 1 else outs[0][1]
                    if len(outs) > 1:
                        st.download_button("Download ZIP", z, file_name=f"converted-{target}.zip", mime="application/zip")
                    else:
                        st.download_button("Download", z, file_name=outs[0][0], mime=f"image/{target}")

        with st.container(border=True):
            st.subheader("Watermark IMAGE")
            st.caption("Stamp text watermark with adjustable opacity.")
            f = st.file_uploader(
                "Image", type=["jpg", "jpeg", "png", "webp", "heic", "heif", "tif", "tiff", "bmp", "gif"], key="img_wm"
            )
            text = st.text_input("Watermark text", value="WATERMARK", key="img_wm_text")
            opacity = st.slider("Opacity", 5, 100, 25, key="img_wm_op")
            if f and st.button("Apply watermark", key="img_wm_btn"):
                from PIL import ImageDraw, ImageFont

                img = pil_open_image(f).convert("RGBA")
                overlay = img.copy()
                draw = ImageDraw.Draw(overlay)
                font = ImageFont.load_default()
                w, h = img.size
                draw.text((w * 0.05, h * 0.9), text, fill=(255, 255, 255, int(255 * opacity / 100)), font=font)
                out = overlay
                data = pil_to_bytes(out, "png")
                st.download_button("Download", data, file_name=f"{Path(f.name).stem}-watermarked.png", mime="image/png")

    st.subheader("Editors")
    edit1, edit2, edit3 = st.columns(3)

    with edit1:
        with st.container(border=True):
            st.subheader("Photo editor")
            st.caption("Adjust brightness/contrast/saturation/sharpness.")
            f = st.file_uploader(
                "Image", type=["jpg", "jpeg", "png", "webp", "heic", "heif", "tif", "tiff", "bmp", "gif"], key="img_edit"
            )
            if f:
                from PIL import ImageEnhance

                img = pil_open_image(f)
                st.image(img, use_container_width=True)
                bright = st.slider("Brightness", 0.2, 2.0, 1.0, 0.05, key="img_edit_b")
                contrast = st.slider("Contrast", 0.2, 2.0, 1.0, 0.05, key="img_edit_c")
                color = st.slider("Saturation", 0.0, 2.0, 1.0, 0.05, key="img_edit_s")
                sharp = st.slider("Sharpness", 0.0, 3.0, 1.0, 0.05, key="img_edit_sh")
                if st.button("Apply edits", key="img_edit_btn"):
                    out = ImageEnhance.Brightness(img).enhance(bright)
                    out = ImageEnhance.Contrast(out).enhance(contrast)
                    out = ImageEnhance.Color(out).enhance(color)
                    out = ImageEnhance.Sharpness(out).enhance(sharp)
                    data = pil_to_bytes(out, "png")
                    st.download_button(
                        "Download",
                        data,
                        file_name=f"{Path(f.name).stem}-edited.png",
                        mime="image/png",
                    )

    with edit2:
        with st.container(border=True):
            st.subheader("Meme generator")
            st.caption("Add top/bottom text.")
            f = st.file_uploader(
                "Image", type=["jpg", "jpeg", "png", "webp", "heic", "heif", "tif", "tiff", "bmp", "gif"], key="img_meme"
            )
            top = st.text_input("Top text", key="img_meme_top")
            bottom = st.text_input("Bottom text", key="img_meme_bottom")
            if f and st.button("Generate", key="img_meme_btn"):
                from PIL import ImageDraw, ImageFont

                img = pil_open_image(f).convert("RGB")
                draw = ImageDraw.Draw(img)
                font = ImageFont.load_default()
                w, h = img.size

                def outline_text(x, y, t):
                    for dx in (-1, 0, 1):
                        for dy in (-1, 0, 1):
                            draw.text((x + dx, y + dy), t, font=font, fill=(0, 0, 0))
                    draw.text((x, y), t, font=font, fill=(255, 255, 255))

                if top:
                    outline_text(10, 10, top)
                if bottom:
                    outline_text(10, h - 25, bottom)
                data = pil_to_bytes(img, "png")
                st.download_button(
                    "Download",
                    data,
                    file_name=f"{Path(f.name).stem}-meme.png",
                    mime="image/png",
                )

    with edit3:
        with st.container(border=True):
            st.subheader("Blur face")
            st.caption("Auto face blur uses OpenCV if installed; otherwise use manual boxes.")
            f = st.file_uploader(
                "Image", type=["jpg", "jpeg", "png", "webp", "heic", "heif", "tif", "tiff", "bmp", "gif"], key="img_blur"
            )
            radius = st.slider("Blur radius", 3, 40, 12, key="img_blur_r")
            auto = st.checkbox("Auto-detect faces (requires opencv-python)", value=False, key="img_blur_auto")
            st.caption("Manual boxes format: x1,y1,x2,y2 per line in pixels")
            boxes_text = st.text_area("Manual boxes", height=90, key="img_blur_boxes")
            if f and st.button("Blur", key="img_blur_btn"):
                img = pil_open_image(f)
                boxes: list[tuple[int, int, int, int]] = []
                if auto:
                    try:
                        import cv2  # type: ignore
                        import numpy as np  # type: ignore

                        rgb = img.convert("RGB")
                        arr = np.array(rgb)
                        gray = cv2.cvtColor(arr, cv2.COLOR_RGB2GRAY)
                        cascade = cv2.CascadeClassifier(cv2.data.haarcascades + "haarcascade_frontalface_default.xml")
                        faces = cascade.detectMultiScale(gray, 1.1, 5)
                        for (x, y, w, h) in faces:
                            boxes.append((int(x), int(y), int(x + w), int(y + h)))
                    except Exception as e:
                        st.warning(f"Auto-detect unavailable: {e}")

                for line in (boxes_text or "").splitlines():
                    line = line.strip()
                    if not line:
                        continue
                    parts = [p.strip() for p in line.split(",")]
                    if len(parts) != 4:
                        st.error(f"Invalid box line: {line}")
                        st.stop()
                    boxes.append((int(parts[0]), int(parts[1]), int(parts[2]), int(parts[3])))

                if not boxes:
                    st.warning("No boxes found (auto or manual).")
                    st.stop()
                out = blur_boxes(img, boxes, radius=int(radius))
                data = pil_to_bytes(out, "png")
                st.download_button(
                    "Download",
                    data,
                    file_name=f"{Path(f.name).stem}-blurred.png",
                    mime="image/png",
                )



if __name__ == "__main__":
    main()
