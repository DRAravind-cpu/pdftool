import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

import streamlit as st
from app import *

st.set_page_config(page_title=f"{APP_TITLE} - PDF to POWERPOINT", layout="wide")
if st.button("← Home"):
    st.switch_page("app.py")

st.title("PDF to POWERPOINT")

deps = get_dependency_status()

st.info("Exports each page as an image slide (requires Poppler).")
if not deps.poppler:
    st.warning("Requires 'pdftoppm' (poppler).")
f = st.file_uploader("PDF", type=["pdf"], key="pdf2ppt")
if f and st.button("Convert", key="pdf2ppt_btn"):
    if not deps.poppler:
        st.stop()
    from pptx import Presentation
    from pptx.util import Inches

    images = pdf2image_convert_from_bytes(f.read(), fmt="png", poppler_path=deps.poppler_path)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for img in images:
        slide = prs.slides.add_slide(blank)
        img_buf = io.BytesIO()
        img.save(img_buf, format="PNG")
        img_buf.seek(0)
        slide.shapes.add_picture(img_buf, Inches(0), Inches(0), width=prs.slide_width)
    out = io.BytesIO()
    prs.save(out)
    st.download_button(
        "Download PPTX",
        out.getvalue(),
        file_name="slides.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

